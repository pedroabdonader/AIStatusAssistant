from flask import Flask, request, send_file, render_template, jsonify
import os
import json
import requests
from openai import AzureOpenAI
import pandas as pd
import ast
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from io import BytesIO
import pytz

app = Flask(__name__)

# Azure OpenAI client configuration
endpoint = os.environ.get("AZURE_ENDPOINT")
model_name = "gpt-4o-mini"
deployment = "gpt-4o-mini"

subscription_key = os.environ.get("AZURE_KEY")
api_version = "2025-03-01-preview"

client = AzureOpenAI(
    api_version=api_version,
    azure_endpoint=endpoint,
    api_key=subscription_key,
)

def get_response(user_input):
    messages = [
        {
            "role": "system",
            "content": (
                "You are a status summarization assistant designed to process project updates. "
                "Your task is to summarize the status of various workstreams based on the provided notes. "
                "You should only respond based on the provided notes. If the notes are insufficient or irrelevant, indicate that there is not enough information to summarize the status update."
                "For each workstream, you should extract and report the following information: "
                "1. Workstream (name of the workstream) "
                "2. Status (choose from: Done, On Time, At Risk, Late, Cancelled) "
                "3. Achievements (a brief summary of weekly achievements, including completion dates in MM/DD format if available) "
                "4. Next Steps (a brief summary of upcoming tasks, including expected completion dates in MM/DD format if available) "
                "5. Expected End Date (if available, the anticipated completion date in MM/DD format). "
                "Additionally, provide a summary for all workstreams that includes: "
                "1. Title (a concise title with key information and the date in MM/DD format) "
                "2. Description (a summary of all status updates) "
                "3. Key Decisions (important decisions made during the update period) "
                "4. Issues/Risks (any identified issues or risks)."
            )
        },
        {
            "role": "user",
            "content": """Here are my notes as of today 2025-05-07: Technical Scrum Call 5/7/2025 
            BPM: 
            A. Blockers:  
            B. What did you do this week:  Looked into SBX data, found some areas missing data (will follow-up on stand up call 5/9) – emailed BPO have not heard back yet so there is data for testing. Revamp BRD and process flows to ensure tech requirements are captured. SNOW team requirements to start onboarding process awaiting to hear latest update from Developer. 
            C. What are the goals for the week:  Clean-up test data and review requirements on business stand-up call. Please include IT Lead on meeting (5/14) with BPO to include AI conversation as well. 
            
            DE: 
            A. Blockers: n/a 
            B. What did you do this week:  QDR documents completed – UAT performed last week. DE to push the QDR documents to ODL. Impact to timeline is related to the table (managing of the repository) continuing to work on timeline impact is TBD.  Team to follow-up with DE lead on 5/7 will connect with PM. IOQ draft to be ready on 5/12. 
            C. What are the goals for the week: DE team to finish managing tasks on Datalake and perform UAT on all tables that were ingestion then can proceed to validation. team to follow-up on validated data timeline. 
            
            AI: 
            A. Blockers: AI developer not assigned to SOP. 
            B. What did you do this week: IT leadership are working with AI leadership to provide POC for AI developer. Expecting AI developer by EOD, 5/7. 
            C. What are the goals for the week: Assign AI developer to SOP. Need to uderstand AI dev timeline working along closely with SA. 
            
            SA: 
            A. Blockers: Still awaiting AI developer 
            B. What did you do this week: n/a 
            C. What are the goals for the week: n/a 
            
            PM/IL: 
            A. Blockers: Complete SRA awaiting AI developer to finalize the deliverable for review 
            B. What did you do this week: Re-working with the PP after the AI 
            C. What are the goals for the week: Update PP with AI timelines once SA questions and resource is confirmed. Need to determine impact of AI development to timeline around integrations. Workstream leads to update status report by 5/8, 12pm. """
        },
        {
            "role": "assistant",
            "content": """{
                "updates": [
                    {
                    "workstream": "BPM",
                    "accomplishments": [
                        "Appian UI has been developed",
                        "Awaiting AI team member",
                        "SNOW On-Boarding Workflow Requirements Defined (85%)",
                        "Investigated SBX data and continued to refine BRD with requirements."
                    ],
                    "next_steps": [
                        "Work with the AI team member and showcase the Appian UI/Workflow to date and prepare for Integration",
                        "Walkthrough AI requirements with BPO"
                    ],
                    "expected_end_date": ""
                    },
                    {
                    "workstream": "DE",
                    "accomplishments": [
                        "VAL data ingestion started this week 5/7",
                        "UAT completed on 5/12"
                    ],
                    "next_steps": [
                        "IoQ Draft to be ready on 5/17",
                        "Finish managing tasks on datalake and perform UAT on all tables"
                    ],
                    "expected_end_date": ""
                    },
                    {
                    "workstream": "AI",
                    "accomplishments": [
                        "Developer not assigned yet",
                        "IT leadership met with AI leadership to provide POC for developer"
                    ],
                    "next_steps": [
                        "SG1 completed, TRA is WIP, SG2 and SG3 requires input from AI Developer for finalizing technical stack",
                        "Understand AI dev timeline"
                    ],
                    "expected_end_date": ""
                    }
                ],
                "issues_risks": [
                    "Need AI developer for search component, SDER SG2 review, and System Risk Assessment (SRA)",
                    "BPM workstream on hold until AI development is at least 50% done to begin integration testing"
                ],
                "key_decisions": [
                    "Estimated Go live date moved to 6/26",
                    "AI dev timeline estimated 6-7 weeks",
                    "SRA to be signed by BPO"
                ],
                "description": "This update summarizes the current status of various workstreams including BPM, DE, AI, SA, and PM/IL as of May 20, 2025. Key achievements include completion of QDR documents and ongoing efforts to onboard an AI developer. Several workstreams are facing delays due to the lack of an assigned AI developer, impacting timelines and deliverables.",
                "title": "Weekly Status Update - 05/20"
                }"""
        },
        {"role": "user", "content": f"Here are my notes as of today {datetime.now(tz=pytz.timezone('America/New_York')).strftime('%Y-%m-%d')}: {user_input}"}
    ]

    response_format = {
    "format": {
        "type": "json_schema",
        "name": "status_summarization",
        "schema": {
            "type": "object",
            "properties": {
            "updates": {
                "type": "array",
                "description": "A list of workstreams with their respective status information.",
                "items": {
                "type": "object",
                "properties": {
                    "Workstream": {
                    "type": "string",
                    "description": "The name of the workstream."
                    },
                    "Status": {
                    "type": "string",
                    "description": "The current status of the workstream.",
                    "enum": [
                        "Done",
                        "On Time",
                        "At Risk",
                        "Late",
                        "Cancelled"
                    ]
                    },
                    "Achievements": {
                    "type": "string",
                    "description": "A short summary of the weekly achievements with completion dates in MM/DD format if available."
                    },
                    "Next Steps": {
                    "type": "string",
                    "description": "A short summary of the next steps with expected completion dates in MM/DD format if available."
                    },
                    "Expected End Date": {
                    "type": "string",
                    "description": "The expected completion date in MM/DD format for this specific workstream."
                    }
                },
                "required": [
                    "Workstream",
                    "Status",
                    "Achievements",
                    "Next Steps",
                    "Expected End Date"
                ],
                "additionalProperties": False
                }
            },
            "title": {
                "type": "string",
                "description": "The title summarizing key information across workstreams with the date."
            },
            "description": {
                "type": "string",
                "description": "A summary description of the status updates across all workstreams."
            },
            "key_decisions": {
                "type": "array",
                "description": "Key decisions made during the status update period.",
                "items": {
                    "type": "string",
                    "description": "A description of the key decision."
                }
            },
            "issues_risks": {
                "type": "array",
                "description": "Issues and risks identified in the update.",
                "items": {
                    "type": "string",
                    "description": "A description of the issue or risk."
                }
            },
            "enough_information": {
                "type": "string",
                "description": "Indicates whether there is enough information to summarize the status update.",
                "enum": [
                    "True",
                    "False"
                ]
            }
            },
            "required": [
            "updates",
            "title",
            "description",
            "key_decisions",
            "issues_risks",
            "enough_information"
            ],
            "additionalProperties": False
        },
        "strict": True
        }
    }





    response = client.responses.create(
        stream=False,
        input=messages,
        text=response_format,
        temperature=0,
        top_p=1.0,
        model=deployment,
    )
    print(response.output_text)  # Debugging line
    response_dict = ast.literal_eval(response.output_text)

    print(response_dict)  # Debugging line
    print(type(response_dict))  # Debugging line
    return response_dict

def createDf(data):
    df = pd.DataFrame(data)
    return df

def populate_powerpoint_template(df, title, description, key_decisions, issues_risks):
    presentation = Presentation("template.pptx")  # Ensure you have a template.pptx file

    #Get the cover slide
    slide = presentation.slides[0]

    for shape in slide.shapes:
        if shape.name == "Cover Title":
            shape.text = title
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)


    # Get the Status slide
    slide = presentation.slides[1]

    # Define the position and size of the table
    left = Inches(0.53)
    top = Inches(2.25)
    # Define column widths
    col_widths = [Inches(1.14), Inches(0.8), Inches(2.64), Inches(2.64), Inches(1.79)]
    rows, cols = df.shape

    # Add a table to the slide
    table = slide.shapes.add_table(rows + 1, cols, left, top, sum(col_widths), Inches(0.5 + rows * 0.5)).table

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Set header row height
    table.rows[0].height = Inches(0.3)  # Set header row height to 0.3 inches

    # Set header row with background color and font size
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 250, 204)  # #FFFACC
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Populate the table with DataFrame data and format other rows
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(df.iat[row_idx, col_idx])

            if isinstance(df.iat[row_idx, col_idx],list):   # Check if the cell contains a list, if it does, join the list into a string in different lines (This is for the Achievements and Next Steps columns)
                cell.text = str("\n".join(df.iat[row_idx, col_idx]))

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 250, 250)  #White background
            # Conditional formatting for the second column (status)
            if col_idx == 1:  # Check if it's the second column
                status = str(df.iat[row_idx, col_idx])
                if status == "On Time":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(210, 244, 220)  # Light Green
                elif status == "At Risk":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 240, 204)  # Light Yellow
                elif status == "Late":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 217, 215)  # Light Red
                elif status == "Cancelled":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(191, 191, 191)  # Light Gray
                elif status == "Done":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(208, 232, 250)  # Light Blue
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Default to white background
            # Set font size and color for all paragraphs in the cell
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black font


    # Populate other shapes with the title, description, key decisions, and issues/risks
    # Format Key Decisions Box with bullets
    for shape in slide.shapes:
        if shape.name == "Description":
            shape.text = description
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)    
        elif shape.name == "Title":
            shape.text = title
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
        elif shape.name == "Key Decisions Box":
            shape.text = "\n".join(key_decisions)
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
        elif shape.name == "Issues/Risks Box":
            shape.text = "\n".join(issues_risks)
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)



    # Save the presentation to a BytesIO object
    ppt_stream = BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)  # Move to the beginning of the stream
    return ppt_stream


def set_shape_format(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0, 0, 0)

def set_shape_fill_color(shape, status):
    color_mapping = {
        "On Time": RGBColor(0, 176, 80),
        "At Risk": RGBColor(255, 192, 0),
        "Late": RGBColor(192, 0, 0),
        "Done": RGBColor(0, 112, 192),
        "Cancelled": RGBColor(125, 125, 125)
    }
    fill_color = color_mapping.get(status, RGBColor(255, 255, 255))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        notes = request.form['notes']
        result = get_response(notes)
        print(result)  # Debugging line
        if result['enough_information'] == "False":
            return jsonify({"columns": ["Error"], "data":[{"Error":"Insufficient Notes for summary generation, please provide more notes and try again."}]}),400
        else:
            df = createDf(result['updates'])
            
            # Include column names in the response
            response_data = {
                'columns': df.columns.tolist(),  # Get the column names
                'data': df.to_dict(orient='records')  # Get the data
            }
            print(result['title'])
            print(result['description'])
            print(result['key_decisions'])
            print(result['issues_risks'])

            print(type(result['key_decisions']))
            print(type(result['issues_risks']))

            # Save the PowerPoint file to a BytesIO object
            ppt_stream = populate_powerpoint_template(df, result['title'], result['description'], result['key_decisions'], result['issues_risks'])

            # Store the stream in a global variable or use a session to access it later
            global ppt_file_stream
            ppt_file_stream = ppt_stream

            print('-----------------------------------------------------')
            print(response_data)

            # Return the DataFrame as JSON
            return jsonify(response_data)

    return render_template('index.html')

@app.route('/download', methods=['GET'])
def download():
    # Send the PowerPoint file for download
    return send_file(ppt_file_stream, as_attachment=True, download_name=f"Status_Report_{datetime.now(tz=pytz.timezone('America/New_York')).strftime('%Y-%m-%d')}.pptx")


if __name__ == '__main__':
    app.run(debug=True)
